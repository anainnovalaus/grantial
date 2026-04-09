from openai import OpenAI
from docx import Document
from pptx import Presentation
import pdfplumber
from utils.postgreSQL import get_connection
import json
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# PDF
def extract_text_from_pdf(path):
    text = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text.append(page.extract_text() or "")
    return "\n".join(text)

# DOCX
def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

# PPTX
def extract_text_from_pptx(path):
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(shape.text)
    return "\n".join(slides_text)

# Función para extraer texto según el tipo de archivo
def extract_text(path, file_extension):
    if file_extension == ".pdf":
        print("El documento es un PDF, extrayendo texto...")
        return extract_text_from_pdf(path)
    elif file_extension == ".docx":
        print("El documento es un DOCX, extrayendo texto...")
        return extract_text_from_docx(path)
    elif file_extension == ".pptx":
        print("El documento es un PPTX, extrayendo texto...")
        return extract_text_from_pptx(path)
    else:
        raise ValueError(f"Formato no soportado: {file_extension}")
    
# Recuperar el texto del fichero
def get_entity_analysis_results(entity_id):
    """
    Recupera todos los 'analysis_result' de la tabla entity_documents
    para el entity_id especificado y los devuelve en una lista.
    
    :param entity_id: ID de la entidad a consultar
    :return: List[str] con todos los valores de analysis_result
    :raises: Exception en caso de error de conexión o de ejecución de la query
    """
    logger.info(f"Recuperando resultados de análisis para entity_id={entity_id}...")
    conn = None
    try:
        # 1. Conectarse
        conn = get_connection()
        cursor = conn.cursor()
        
        # 2. Ejecutar la consulta
        query = """
            SELECT analysis_result
            FROM entity_documents
            WHERE entity_id = %s
        """
        cursor.execute(query, (entity_id,))
        
        # 3. Leer todos los rows
        rows = cursor.fetchall()  # devuelve lista de tuplas
        
        # 4. Extraer columna en una lista plana
        results = [r[0] for r in rows]
        logger.info(f"Recuperados {len(results)} resultados de análisis para entity_id={entity_id}.")
        return "\n".join(results)

    except Exception as e:
        logger.error(f"Error al recuperar analysis_result para entity_id={entity_id}: {e}")
        raise

    finally:
        # 5. Cerrar cursor y conexión
        if cursor:
            cursor.close()
        if conn:
            conn.close()

# Chatgpt analiza la información de la entidad
def analyze_with_gpt(client, pagina_web, razon_social, nif, descripcion_usuario, document_text, OPENAI_PROMPT_SCRAPE_ENTITY):

    if pagina_web: 
        response = client.responses.create(
            model="gpt-4.1-mini",
            tools=[{"type": "web_search_preview"}],
            input=f"""Eres un asistente que extrae toda la información posible de una página web.
                Recibirás el link a la página y debes devolver un JSON estructurado
                con campos como título, meta tags, encabezados, párrafos, enlaces, imágenes, tablas, JSON-LD y datos de contacto (telefono y email).
                En concreto, debes analizar la entidad {razon_social} con CIF {nif} y descripción: {descripcion_usuario} y el siguiente link: {pagina_web}"""
        )
        response = response.output_text
        logger.info(response)
    else:
        response = "No se ha proporcionado página web."
        logger.info(response)

    response2 = client.responses.create(
        model="gpt-4.1-mini",
        tools=[{"type": "web_search_preview"}],
        input=f"""Eres un asistente que busca toda la información posible de una entidad privada española.
            Recibirás el nombre social y el CIF de la entidad y debes devolver texto descriptivo
            con información relevante como actividad, dirección, teléfono, email, objeto social, fecha de constitución
            sector, capital social, administradores, cnae, y otros datos relevantes.
            En concreto, debes analizar la entidad {razon_social} con CIF {nif} y descripción: {descripcion_usuario} """
    )

    final_response = client.responses.create(
        model="gpt-4.1-nano",
        input=[
            {
                "role": "developer",
                "content": OPENAI_PROMPT_SCRAPE_ENTITY
            },
            {
                "role": "user",
                "content": "Proporciono la siguiente información extraida de internet:\n" + response2.output_text + "\n\n" + response +
                "\n\n Información extraida de los documentos:" + document_text
            }
        ]
    )

    print("=== JSON generado ===")
    print(final_response.output_text)
    return final_response.output_text

"""===================="""
# MAIN FUNCTION
def scrape_main(pagina_web, razon_social, nif, entity_id, descripcion_usuario, uploaded_files, OPENAI_PROMPT_SCRAPE_ENTITY):

    client = OpenAI()

    try:
        logger.info(f"=== Iniciando el análisis del documento de la entidad {razon_social}===")

        if uploaded_files:
            logger.info("Se han encontrado documentos para analizar.")
            document_text = get_entity_analysis_results(entity_id)
        else:
            document_text = "No se han proporcionado documentos para analizar."
            logger.info("No se han proporcionado documentos para analizar.")

        raw = analyze_with_gpt(client, pagina_web, razon_social, nif, descripcion_usuario, document_text, OPENAI_PROMPT_SCRAPE_ENTITY)

        if raw:
            logger.info("=== Análisis completado ===")
            # Parseamos la cadena JSON a dict
            try:
                final_response = json.loads(raw)
                logger.info(f"Parsed scraping JSON (type={type(final_response)}): {final_response}")
            except json.JSONDecodeError as err:
                logger.error("Error al parsear JSON de scraping:", err)
                return
            return final_response

        else:
            logger.error("No se pudo completar el análisis.")
            return 
        
    except Exception as e:
        logger.error(f"Error al procesar el fichero: {e}")
        return
